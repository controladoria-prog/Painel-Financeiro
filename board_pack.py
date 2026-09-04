# -*- coding: utf-8 -*-
"""Board pack -- a apresentação do fechamento mensal, pronta no dia 5.

Gera um PPTX de seis slides (capa, semáforo e cartões, receita por mês,
EBITDA por mês, estouros e folgas, narrativa) com os MESMOS números e a
MESMA narrativa da tela e do e-mail, e manda por e-mail como anexo. A
reunião de resultados começa com o material pronto -- e igual todo mês.

Uso:
    python board_pack.py                    gera e envia
    python board_pack.py --teste            gera e grava board_pack.pptx, sem enviar
    python board_pack.py --salvar arq.pptx  grava com outro nome

Precisa de: pip install python-pptx (além de pandas, numpy, openpyxl).
"""
import io
import os
import sys

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt

from briefing import (DIAS_SEMANA, carregar_funcoes_do_app, enviar_email, montar_briefing,
                      series_mensais, status_geral)

NAVY = RGBColor(0x1B, 0x2A, 0x41)
CINZA = RGBColor(0x6B, 0x72, 0x80)
TEXTO = RGBColor(0x1F, 0x29, 0x37)
VERDE = RGBColor(0x1E, 0x84, 0x49)
VERMELHO = RGBColor(0xC0, 0x39, 0x2B)
AMBAR = RGBColor(0xB9, 0x77, 0x0E)
AZUL = RGBColor(0x2F, 0x6F, 0xD6)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
TONS = {"negativo": VERMELHO, "positivo": VERDE, "alerta": AMBAR, "neutro": CINZA}


def _texto(slide, x, y, w, h, texto, tamanho=14, cor=TEXTO, negrito=False, alinhamento=None):
    caixa = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    quadro = caixa.text_frame
    quadro.word_wrap = True
    paragrafo = quadro.paragraphs[0]
    paragrafo.text = str(texto)
    paragrafo.font.size = Pt(tamanho)
    paragrafo.font.bold = negrito
    paragrafo.font.color.rgb = cor
    if alinhamento is not None:
        paragrafo.alignment = alinhamento
    return caixa


def _retangulo(slide, x, y, w, h, cor):
    from pptx.enum.shapes import MSO_SHAPE
    forma = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    forma.fill.solid()
    forma.fill.fore_color.rgb = cor
    forma.line.fill.background()
    return forma


def _cabecalho(slide, titulo, subtitulo=""):
    _retangulo(slide, 0, 0, 13.333, 0.9, NAVY)
    _texto(slide, 0.5, 0.18, 9, 0.5, titulo, 22, BRANCO, True)
    if subtitulo:
        _texto(slide, 9.3, 0.3, 3.6, 0.4, subtitulo, 11, RGBColor(0x9F, 0xB3, 0xD1))


def _grafico_colunas(slide, titulo, rotulos, real, orc, y):
    dados = CategoryChartData()
    dados.categories = rotulos
    dados.add_series("Realizado", [v / 1e6 for v in real])
    dados.add_series("Orçado", [v / 1e6 for v in orc])
    grafico = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(y),
                                     Inches(12.3), Inches(5.3), dados).chart
    grafico.has_legend = True
    grafico.legend.position = XL_LEGEND_POSITION.BOTTOM
    grafico.legend.include_in_layout = False
    grafico.value_axis.has_major_gridlines = False
    grafico.value_axis.minimum_scale = 0   # eixo do zero: recorte inflava a diferença
    grafico.value_axis.tick_labels.font.size = Pt(10)
    grafico.category_axis.tick_labels.font.size = Pt(11)
    grafico.plots[0].series[0].format.fill.solid()
    grafico.plots[0].series[0].format.fill.fore_color.rgb = AZUL
    grafico.plots[0].series[1].format.fill.solid()
    grafico.plots[0].series[1].format.fill.fore_color.rgb = RGBColor(0xC9, 0xD1, 0xDB)
    grafico.plots[0].has_data_labels = True
    grafico.plots[0].data_labels.number_format = '0.0"M"'
    grafico.plots[0].data_labels.font.size = Pt(9)
    return grafico


def montar_board_pack(fatos, itens, series, hoje, ns=None, modo="consolidado"):
    """Devolve os bytes do PPTX. `modo="departamento"` troca receita/EBITDA
    por gasto realizado x orçado do departamento (fatos_do_departamento)."""
    if modo == "departamento":
        return _board_pack_departamento(fatos, itens, series, hoje, ns)
    fmt = (ns or {}).get("formata_valor_curto", lambda v: f"R$ {v:,.0f}")
    pct = (ns or {}).get("_pct_br", lambda v, casas=1: f"{v:.{casas}f}".replace(".", ","))
    apresentacao = Presentation()
    apresentacao.slide_width = Inches(13.333)
    apresentacao.slide_height = Inches(7.5)
    em_branco = apresentacao.slide_layouts[6]
    dia = f"{DIAS_SEMANA[hoje.weekday()]}, {hoje.strftime('%d/%m/%Y')}"
    status, _cor_status, frase_status = status_geral(fatos)

    # 1) capa
    slide = apresentacao.slides.add_slide(em_branco)
    _retangulo(slide, 0, 0, 13.333, 7.5, NAVY)
    _texto(slide, 0.8, 2.3, 11, 0.5, "CONTROLADORIA B&A · FECHAMENTO", 13, RGBColor(0x9F, 0xB3, 0xD1))
    _texto(slide, 0.8, 2.8, 11, 1.2, fatos.get("periodo", ""), 40, BRANCO, True)
    _texto(slide, 0.8, 4.1, 11, 0.6, f"{frase_status} · {dia}" if frase_status else dia, 16,
           RGBColor(0x9F, 0xB3, 0xD1))

    # 2) semáforo e cartões
    slide = apresentacao.slides.add_slide(em_branco)
    _cabecalho(slide, "Resultado do período", dia)
    cor_status = TONS.get({"NO ORÇADO": "positivo", "OBSERVAR": "alerta"}.get(status, "negativo"))
    _retangulo(slide, 0.5, 1.2, 12.3, 0.6, cor_status)
    _texto(slide, 0.7, 1.3, 12, 0.4, f"{status} · {frase_status}", 14, BRANCO, True)
    cartoes = []
    if fatos.get("rec_real") is not None:
        var = (fatos["rec_real"] / fatos["rec_orc"] - 1) * 100 if fatos.get("rec_orc") else None
        cartoes.append(("Receita líquida", fmt(fatos["rec_real"]),
                        f"{'▲' if var >= 0 else '▼'} {pct(abs(var))}% vs orçado" if var is not None else ""))
    if fatos.get("ebitda_real") is not None:
        var = (fatos["ebitda_real"] / fatos["ebitda_orc"] - 1) * 100 if fatos.get("ebitda_orc") else None
        cartoes.append(("EBITDA", fmt(fatos["ebitda_real"]),
                        f"{'▲' if var >= 0 else '▼'} {pct(abs(var))}% vs orçado" if var is not None else ""))
    if fatos.get("rec_real"):
        margem = fatos["ebitda_real"] / fatos["rec_real"] * 100
        sub = (f"fecha em {pct(fatos['margem_proj'])}% com os lançamentos pendentes"
               if fatos.get("margem_proj") is not None else "realizada no período")
        cartoes.append(("Margem EBITDA", f"{pct(margem)}%", sub))
    r = fatos.get("ritmo")
    if r:
        cartoes.append((f"Ritmo de {str(r['mes']).capitalize()}", f"{r['pct']:.0f}%",
                        f"dia {r['dia']} de {r['dias']}" + (f" · chance {r['chance'] * 100:.0f}%" if r.get("chance") is not None else "")))
    largura = 12.3 / max(len(cartoes), 1)
    for i, (rotulo, valor, sub) in enumerate(cartoes):
        x = 0.5 + i * largura
        _retangulo(slide, x + 0.05, 2.2, largura - 0.1, 0.06, NAVY)
        _texto(slide, x + 0.15, 2.35, largura - 0.3, 0.4, rotulo.upper(), 10, CINZA)
        _texto(slide, x + 0.15, 2.7, largura - 0.3, 0.8, valor, 30, TEXTO, True)
        _texto(slide, x + 0.15, 3.5, largura - 0.3, 0.8, sub, 11, CINZA)
    _texto(slide, 0.5, 4.8, 12.3, 1.2,
           "Base da análise: meses fechados; o mês corrente entra só como explicação do gap. "
           "Lançamentos chegam D+2; a chance de bater a meta é uma estimativa a partir do histórico do ano.",
           11, CINZA)

    # 3) receita por mês
    slide = apresentacao.slides.add_slide(em_branco)
    _cabecalho(slide, "Receita líquida por mês · realizado vs. orçado (R$ milhões)", dia)
    _grafico_colunas(slide, "Receita", series["rotulos"], series["rec_real"], series["rec_orc"], 1.2)

    # 4) EBITDA por mês
    slide = apresentacao.slides.add_slide(em_branco)
    _cabecalho(slide, "EBITDA por mês · realizado vs. orçado (R$ milhões)", dia)
    _grafico_colunas(slide, "EBITDA", series["rotulos"], series["eb_real"], series["eb_orc"], 1.2)

    # 5) estouros e folgas
    slide = apresentacao.slides.add_slide(em_branco)
    _cabecalho(slide, "Onde passou e onde sobrou (meses fechados)", dia)
    linhas_tab = [("Conta", "Situação", "Valor", "%")]
    for e in fatos.get("estouros") or []:
        linhas_tab.append((e["conta"], "acima do orçado", f"+{fmt(e['desvio'])}",
                           f"+{e['pct']:.0f}%" if e.get("pct") is not None else "s/ orç."))
    for c in fatos.get("folgas") or []:
        linhas_tab.append((c["conta"], "folga real", f"−{fmt(c['folga'])}",
                           f"−{c['pct']:.0f}%" if c.get("pct") is not None else ""))
    for a in fatos.get("artefatos") or []:
        linhas_tab.append((a["conta"], "lançamento pendente (não é economia)", f"−{fmt(a['folga'])}",
                           f"{fmt(a['pendente'])} por entrar"))
    tabela = slide.shapes.add_table(len(linhas_tab), 4, Inches(0.5), Inches(1.2), Inches(12.3),
                                    Inches(0.4 * len(linhas_tab))).table
    for i, linha in enumerate(linhas_tab):
        for j, valor in enumerate(linha):
            celula = tabela.cell(i, j)
            celula.text = str(valor)
            paragrafo = celula.text_frame.paragraphs[0]
            paragrafo.font.size = Pt(12 if i else 11)
            paragrafo.font.bold = i == 0
            paragrafo.font.color.rgb = BRANCO if i == 0 else TEXTO
            celula.fill.solid()
            celula.fill.fore_color.rgb = NAVY if i == 0 else (RGBColor(0xF4, 0xF7, 0xFB) if i % 2 else BRANCO)
    tabela.columns[0].width = Inches(5.3)
    tabela.columns[1].width = Inches(3.5)
    tabela.columns[2].width = Inches(1.8)
    tabela.columns[3].width = Inches(1.7)

    # 6) narrativa
    slide = apresentacao.slides.add_slide(em_branco)
    _cabecalho(slide, "O que aconteceu e por quê", dia)
    y = 1.2
    for item in itens:
        texto_limpo = item["texto"].replace("<b>", "").replace("</b>", "")
        _retangulo(slide, 0.5, y + 0.08, 0.06, 0.55, TONS.get(item.get("tom"), CINZA))
        _texto(slide, 0.7, y, 1.6, 0.4, item["rotulo"].upper(), 10, TONS.get(item.get("tom"), CINZA), True)
        _texto(slide, 2.3, y - 0.02, 10.5, 0.95, texto_limpo, 13, TEXTO)
        y += 0.95
    saida = io.BytesIO()
    apresentacao.save(saida)
    return saida.getvalue()


def _board_pack_departamento(fatos, itens, series, hoje, ns=None):
    fmt = (ns or {}).get("formata_valor_curto", lambda v: f"R$ {v:,.0f}")
    pct = (ns or {}).get("_pct_br", lambda v, casas=1: f"{v:.{casas}f}".replace(".", ","))
    apresentacao = Presentation()
    apresentacao.slide_width = Inches(13.333)
    apresentacao.slide_height = Inches(7.5)
    em_branco = apresentacao.slide_layouts[6]
    dia = f"{DIAS_SEMANA[hoje.weekday()]}, {hoje.strftime('%d/%m/%Y')}"
    gasto_real, gasto_orc = fatos.get("gasto_real", 0.0), fatos.get("gasto_orc", 0.0)
    gap = (gasto_real / gasto_orc - 1) if gasto_orc else 0.0
    status = "NO ORÇADO" if gap <= 0 else ("OBSERVAR" if gap <= 0.05 else "ATENÇÃO")
    cor_status = {"NO ORÇADO": VERDE, "OBSERVAR": AMBAR}.get(status, VERMELHO)
    frase = f"Gasto {pct(abs(gap) * 100)}% {'acima' if gap > 0 else 'abaixo'} do orçado" if gasto_orc else ""
    nome = str(fatos.get("departamento", ""))

    slide = apresentacao.slides.add_slide(em_branco)
    _retangulo(slide, 0, 0, 13.333, 7.5, NAVY)
    _texto(slide, 0.8, 2.3, 11, 0.5, "CONTROLADORIA B&A · RELATÓRIO DE CUSTOS", 13, RGBColor(0x9F, 0xB3, 0xD1))
    _texto(slide, 0.8, 2.8, 11.5, 1.2, nome, 34, BRANCO, True)
    _texto(slide, 0.8, 4.1, 11, 0.6, f"{fatos.get('periodo', '')} · {frase} · {dia}", 15, RGBColor(0x9F, 0xB3, 0xD1))

    slide = apresentacao.slides.add_slide(em_branco)
    _cabecalho(slide, f"Gasto do departamento · {fatos.get('periodo', '')}", dia)
    _retangulo(slide, 0.5, 1.2, 12.3, 0.6, cor_status)
    _texto(slide, 0.7, 1.3, 12, 0.4, f"{status} · {frase}", 14, BRANCO, True)
    cartoes = [("Gasto realizado", fmt(gasto_real), "meses fechados"),
               ("Orçado", fmt(gasto_orc), "mesmo período"),
               ("Desvio", f"{'+' if gasto_real >= gasto_orc else '−'}{fmt(abs(gasto_real - gasto_orc))}",
                f"{pct(abs(gap) * 100)}% {'acima' if gap > 0 else 'abaixo'}" if gasto_orc else ""),
               ("Contas acima do orçado", f"{fatos.get('n_acima', 0)} de {fatos.get('n_contas', 0)}", "no escopo do departamento")]
    largura = 12.3 / 4
    for i, (rotulo, valor, sub) in enumerate(cartoes):
        x = 0.5 + i * largura
        _retangulo(slide, x + 0.05, 2.2, largura - 0.1, 0.06, NAVY)
        _texto(slide, x + 0.15, 2.35, largura - 0.3, 0.4, rotulo.upper(), 10, CINZA)
        _texto(slide, x + 0.15, 2.7, largura - 0.3, 0.8, valor, 28, TEXTO, True)
        _texto(slide, x + 0.15, 3.5, largura - 0.3, 0.8, sub, 11, CINZA)
    _texto(slide, 0.5, 4.8, 12.3, 1.0,
           "Base: meses fechados; o mês corrente fica fora. Lançamentos chegam D+2. "
           "Folga de conta pode ser lançamento de fechamento ainda por entrar.", 11, CINZA)

    slide = apresentacao.slides.add_slide(em_branco)
    _cabecalho(slide, "Gasto por mês · realizado vs. orçado (R$ milhões)", dia)
    _grafico_colunas(slide, "Gasto", series["rotulos"], series["gasto_real"], series["gasto_orc"], 1.2)

    slide = apresentacao.slides.add_slide(em_branco)
    _cabecalho(slide, "Contas: onde passou e onde sobrou", dia)
    linhas_tab = [("Conta", "Situação", "Valor", "%")]
    for e in fatos.get("estouros") or []:
        linhas_tab.append((e["conta"], "acima do orçado", f"+{fmt(e['desvio'])}",
                           f"+{e['pct']:.0f}%" if e.get("pct") is not None else "s/ orç."))
    for c in fatos.get("folgas") or []:
        linhas_tab.append((c["conta"], "abaixo do orçado", f"−{fmt(c['folga'])}",
                           f"−{c['pct']:.0f}%" if c.get("pct") is not None else ""))
    tabela = slide.shapes.add_table(len(linhas_tab), 4, Inches(0.5), Inches(1.2), Inches(12.3),
                                    Inches(0.4 * len(linhas_tab))).table
    for i, linha in enumerate(linhas_tab):
        for j, valor in enumerate(linha):
            celula = tabela.cell(i, j)
            celula.text = str(valor)
            paragrafo = celula.text_frame.paragraphs[0]
            paragrafo.font.size = Pt(12 if i else 11)
            paragrafo.font.bold = i == 0
            paragrafo.font.color.rgb = BRANCO if i == 0 else TEXTO
            celula.fill.solid()
            celula.fill.fore_color.rgb = NAVY if i == 0 else (RGBColor(0xF4, 0xF7, 0xFB) if i % 2 else BRANCO)
    for k, w in enumerate((5.3, 3.5, 1.8, 1.7)):
        tabela.columns[k].width = Inches(w)

    slide = apresentacao.slides.add_slide(em_branco)
    _cabecalho(slide, "O que aconteceu e por quê", dia)
    y = 1.2
    for item in itens:
        _retangulo(slide, 0.5, y + 0.08, 0.06, 0.55, TONS.get(item.get("tom"), CINZA))
        _texto(slide, 0.7, y, 1.6, 0.4, item["rotulo"].upper(), 10, TONS.get(item.get("tom"), CINZA), True)
        _texto(slide, 2.3, y - 0.02, 10.5, 0.95, item["texto"].replace("<b>", "").replace("</b>", ""), 13, TEXTO)
        y += 0.95
    saida = io.BytesIO()
    apresentacao.save(saida)
    return saida.getvalue()


def main(argv):
    ns = carregar_funcoes_do_app()
    fatos, itens, ctx = montar_briefing(ns, url_fech=os.environ.get("FECHAMENTO_CSV_URL", ""))
    hoje = ctx["hoje"]
    series = series_mensais(ns, ctx["list_df_real"], ctx["list_df_orc"], ctx["meses_cols"], ctx["m_map"],
                            ate_mes=hoje.month - 1 if hoje.month > 1 else 12)
    pptx_bytes = montar_board_pack(fatos, itens, series, hoje, ns)
    nome = f"board_pack_{hoje.strftime('%Y-%m')}.pptx"
    if "--salvar" in argv:
        nome = argv[argv.index("--salvar") + 1]
    if "--teste" in argv or "--salvar" in argv:
        with open(nome, "wb") as arquivo:
            arquivo.write(pptx_bytes)
        print(f"Board pack gravado em {nome} ({len(pptx_bytes) // 1024} KB)")
        return
    _status, _cor, frase = status_geral(fatos)
    assunto = f"Board pack · {ctx['rotulo']}" + (f" · {frase}" if frase else "")
    texto = (f"Segue o board pack do fechamento ({ctx['rotulo']}), gerado automaticamente pelo painel.\n\n"
             + "\n".join(f"{i['rotulo']}: {i['texto'].replace('<b>', '').replace('</b>', '')}" for i in itens))
    html = ("<p>Segue o board pack do fechamento, gerado automaticamente pelo painel.</p><ul>"
            + "".join(f"<li><b>{i['rotulo']}:</b> {i['texto']}</li>" for i in itens) + "</ul>")
    destinos = enviar_email(assunto, html, texto, str(ns.get("LOGO_BEEA_B64") or ""),
                            anexos=[(nome, pptx_bytes, "application",
                                     "vnd.openxmlformats-officedocument.presentationml.presentation")])
    print(f"Board pack enviado para {', '.join(destinos)}: {assunto}")


if __name__ == "__main__":
    main(sys.argv[1:])
